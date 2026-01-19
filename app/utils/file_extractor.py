"""
Utilidad para extraer texto de diferentes formatos de archivo
"""

import io
import logging
from typing import Optional
from PIL import Image
import pytesseract

# Configurar logger
logger = logging.getLogger(__name__)

class FileExtractor:
    """
    Clase para extraer texto de diferentes tipos de archivos
    """
    
    @staticmethod
    def extract_text_from_pdf(content: bytes) -> str:
        """
        Extrae texto de un archivo PDF
        
        Args:
            content: Contenido del archivo PDF en bytes
            
        Returns:
            Texto extraido del PDF
        """
        try:
            from pypdf import PdfReader
            
            pdf_file = io.BytesIO(content)
            reader = PdfReader(pdf_file)
            
            text_parts = []
            logger.info(f"[DOC] Iniciando extraccion de PDF. Total paginas: {len(reader.pages)}")
            
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Pagina {page_num + 1} ---\n{page_text}")
                        logger.debug(f"[OK] Pagina {page_num + 1} extraida: {len(page_text)} caracteres")
                    else:
                        logger.warning(f"[WARN] Pagina {page_num + 1} vacia o sin texto seleccionable")
                except Exception as e:
                    logger.warning(f"[WARN] Error extrayendo pagina {page_num + 1}: {e}")
                    continue
            
            extracted_text = "\n\n".join(text_parts)
            
            if not extracted_text.strip():
                logger.warning("[WARN] PDF procesado pero no se extrajo texto (posible PDF escaneado)")
                return "[No se pudo extraer texto del documento. Es posible que sea un PDF escaneado o imagen.]"
                
            logger.info(f"[OK] PDF procesado exitosamente: {len(extracted_text)} caracteres extraidos")
            return extracted_text
            
        except Exception as e:
            logger.error(f"[ERR] Error procesando PDF: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(content: bytes) -> str:
        """
        Extrae texto de un archivo Word (.docx)
        
        Args:
            content: Contenido del archivo Word en bytes
            
        Returns:
            Texto extraido del documento
        """
        try:
            from docx import Document
            
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            text_parts = []
            
            # Extraer texto de parrafos
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extraer texto de tablas
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"[OK] Word procesado: {len(doc.paragraphs)} parrafos, {len(extracted_text)} caracteres")
            return extracted_text
            
        except Exception as e:
            logger.error(f"[ERR] Error procesando Word: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_excel(content: bytes) -> str:
        """
        Extrae texto de un archivo Excel (.xlsx, .xls)
        
        Args:
            content: Contenido del archivo Excel en bytes
            
        Returns:
            Texto extraido de las hojas
        """
        try:
            from openpyxl import load_workbook
            
            excel_file = io.BytesIO(content)
            workbook = load_workbook(excel_file, data_only=True)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Hoja: {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    # Filtrar valores None y convertir a string
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        text_parts.append(" | ".join(row_values))
            
            extracted_text = "\n".join(text_parts)
            logger.info(f"[OK] Excel procesado: {len(workbook.sheetnames)} hojas, {len(extracted_text)} caracteres")
            return extracted_text
            
        except Exception as e:
            logger.error(f"[ERR] Error procesando Excel: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(content: bytes) -> str:
        """
        Extrae texto de un archivo PowerPoint (.pptx)
        
        Args:
            content: Contenido del archivo PowerPoint en bytes
            
        Returns:
            Texto extraido de las diapositivas
        """
        try:
            from pptx import Presentation
            
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, start=1):
                text_parts.append(f"--- Diapositiva {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                    
                    # Extraer texto de tablas en diapositivas
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                text_parts.append(row_text)
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"[OK] PowerPoint procesado: {len(presentation.slides)} diapositivas, {len(extracted_text)} caracteres")
            return extracted_text
            
        except Exception as e:
            logger.error(f"[ERR] Error procesando PowerPoint: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_image(content: bytes, filename: str = "") -> str:
        """
        Extrae texto de una imagen usando OCR (pytesseract)
        
        Args:
            content: Contenido de la imagen en bytes
            filename: Nombre del archivo (para logging)
            
        Returns:
            Texto extraido de la imagen
        """
        try:
            image_file = io.BytesIO(content)
            image = Image.open(image_file)
            
            # Convertir a RGB si es necesario
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Extraer texto usando OCR
            extracted_text = pytesseract.image_to_string(image, lang='spa+eng')
            
            if extracted_text.strip():
                logger.info(f"[OK] Imagen procesada ({filename}): {len(extracted_text)} caracteres extraidos")
                return extracted_text
            else:
                logger.warning(f"[WARN] No se encontro texto en la imagen: {filename}")
                return ""
            
        except Exception as e:
            logger.error(f"[ERR] Error procesando imagen {filename}: {e}")
            logger.info("[IDEA] Asegurate de tener Tesseract OCR instalado: https://github.com/tesseract-ocr/tesseract")
            return ""
    
    @staticmethod
    def extract_text(content: bytes, filename: str) -> str:
        """
        Extrae texto de un archivo segun su extension
        
        Args:
            content: Contenido del archivo en bytes
            filename: Nombre del archivo con extension
            
        Returns:
            Texto extraido del archivo
        """
        file_extension = filename.lower().split('.')[-1]
        
        extractors = {
            'pdf': FileExtractor.extract_text_from_pdf,
            'docx': FileExtractor.extract_text_from_docx,
            'doc': FileExtractor.extract_text_from_docx,
            'xlsx': FileExtractor.extract_text_from_excel,
            'xls': FileExtractor.extract_text_from_excel,
            'pptx': FileExtractor.extract_text_from_pptx,
            'ppt': FileExtractor.extract_text_from_pptx,
            'png': lambda c: FileExtractor.extract_text_from_image(c, filename),
            'jpg': lambda c: FileExtractor.extract_text_from_image(c, filename),
            'jpeg': lambda c: FileExtractor.extract_text_from_image(c, filename),
            'txt': lambda c: c.decode('utf-8'),
            'md': lambda c: c.decode('utf-8'),
        }
        
        extractor = extractors.get(file_extension)
        
        if extractor:
            try:
                return extractor(content)
            except Exception as e:
                logger.error(f"[ERR] Error extrayendo texto de {filename}: {e}")
                return ""
        else:
            logger.warning(f"[WARN] Formato no soportado: {file_extension}")
            return ""
